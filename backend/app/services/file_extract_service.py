import os
import re
import tempfile

from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from app.services.ocr_service import extract_text_from_image_bytes


IMAGE_EXTENSIONS = [".jpg", ".jpeg", ".png", ".webp"]
MIN_USEFUL_PAGE_WORDS = 12
MIN_MULTI_PAGE_PDF_WORDS = 200
MIN_MULTI_PAGE_PDF_WORDS_PER_PAGE = 2
DEVANAGARI_RE = re.compile(r"[\u0900-\u097F]")
LATIN_LETTER_RE = re.compile(r"[A-Za-z]")
BROKEN_PDF_TOKEN_RE = re.compile(r"[;:`~{}\\¼%]|[A-Za-z]{4,}[A-Z][a-z]")
HINDI_FILENAME_HINTS = (
    "hindi",
    "jhks",
    "kshitij",
    "kritika",
    "sparsh",
    "sanchayan",
    "vasant",
    "durva",
    "rimjhim",
    # NCERT Grade-Hindi textbook file codes (ghXX prefix series)
    "ghml",   # MahaLekhak — Hindi grammar/reader
    "ghsn",   # Sansaar Naati Kahi
    "ghst",   # Hindi supplementary
    "ghsc",   # Hindi supplementary
    "ghmh",   # MahaBharat / Hindi heritage
    "ihga",   # NCERT integrated Hindi Grade series
    "iehe",   # English/Hindi combined
    "hhss",   # Hindi higher secondary
)

def extract_text_from_txt(file_bytes: bytes) -> str:
    """Decode a plain text upload into normalized UTF-8 text."""
    return file_bytes.decode("utf-8", errors="ignore").strip()


def extract_text_from_pdf(file_bytes: bytes, filename: str = "") -> str:
    """Extract embedded text from all pages of a PDF upload."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        try:
            reader = PdfReader(tmp_path)
            text_parts = []

            for page in reader.pages:
                text_parts.append(page.extract_text() or "")

            page_count = len(reader.pages)
        except Exception:
            text_parts = extract_pdf_text_parts_with_pymupdf(file_bytes)
            page_count = len(text_parts)

        extracted_text = "\n\n".join(text_parts).strip()

        if looks_like_broken_hindi_pdf_text(extracted_text, filename):
            ocr_text = extract_rendered_pdf_text_with_ocr(file_bytes)
            if is_better_hindi_extraction(ocr_text, extracted_text):
                extracted_text = ocr_text

        word_count = count_words(extracted_text)

        if (
            page_count >= 10
            and word_count < MIN_MULTI_PAGE_PDF_WORDS
            and word_count < page_count * MIN_MULTI_PAGE_PDF_WORDS_PER_PAGE
        ):
            raise ValueError(
                "This PDF appears to be scanned or image-based. "
                f"Only {word_count} readable words were extracted from {page_count} pages, "
                "so RAG upload was stopped to avoid creating an unusable one-chunk document. "
                "Please upload chapter images, a searchable/OCR PDF, or use the full-book OCR splitter flow."
            )

        return extracted_text

    finally:
        os.remove(tmp_path)


def extract_pdf_text_parts_with_pymupdf(file_bytes: bytes) -> list[str]:
    """Fallback PDF text extraction for compressed PDFs that pypdf cannot parse."""
    try:
        import fitz  # PyMuPDF
    except Exception as exc:
        raise RuntimeError(
            "PDF text extraction failed and PyMuPDF fallback is unavailable. "
            "Install backend requirement PyMuPDF and redeploy."
        ) from exc

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    document = None
    try:
        document = fitz.open(tmp_path)
        return [
            document.load_page(page_index).get_text("text") or ""
            for page_index in range(document.page_count)
        ]
    finally:
        if document is not None:
            document.close()
        os.remove(tmp_path)


def has_hindi_filename_hint(filename: str) -> bool:
    """Return true when a filename strongly suggests Hindi textbook content."""
    lower_filename = (filename or "").lower()
    return any(hint in lower_filename for hint in HINDI_FILENAME_HINTS)


def count_devanagari_chars(text: str) -> int:
    """Count Hindi/Devanagari characters in extracted text."""
    return len(DEVANAGARI_RE.findall(text or ""))


def looks_like_broken_hindi_pdf_text(text: str, filename: str = "") -> bool:
    """
    Detect Hindi PDFs whose embedded text layer decoded as font-encoded noise.

    Some NCERT Hindi PDFs expose a Latin-looking text layer even though the page
    visibly contains Devanagari. Indexing that text would poison RAG, so these
    documents need rendered-page OCR instead of normal PDF text extraction.
    """
    sample = (text or "").strip()[:2500]
    if not sample or not has_hindi_filename_hint(filename):
        return False

    if count_devanagari_chars(sample) >= 20:
        return False

    tokens = re.findall(r"\S+", sample)
    if not tokens:
        return False

    latin_count = len(LATIN_LETTER_RE.findall(sample))
    noisy_token_count = sum(
        1
        for token in tokens
        if BROKEN_PDF_TOKEN_RE.search(token)
        or (
            len(token) >= 5
            and LATIN_LETTER_RE.search(token)
            and not re.search(r"[aeiouAEIOU]", token)
        )
    )
    noisy_ratio = noisy_token_count / max(len(tokens), 1)

    return latin_count >= 80 and noisy_ratio >= 0.18


def is_better_hindi_extraction(candidate_text: str, original_text: str) -> bool:
    """Prefer OCR text only when it clearly improves Devanagari extraction."""
    candidate = (candidate_text or "").strip()
    if not candidate:
        return False

    return (
        count_devanagari_chars(candidate) > max(20, count_devanagari_chars(original_text))
        and count_words(candidate) >= max(10, count_words(original_text) // 4)
    )


def render_pdf_page_to_jpeg(file_bytes: bytes, page_index: int) -> bytes:
    """Render one PDF page into JPEG bytes for OCR fallback."""
    try:
        import fitz  # PyMuPDF
    except Exception as exc:
        raise RuntimeError(
            "Hindi PDF text extraction needs PyMuPDF for rendered-page OCR. "
            "Install backend requirement PyMuPDF and redeploy."
        ) from exc

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    document = None
    try:
        document = fitz.open(tmp_path)
        page = document.load_page(page_index)
        matrix = fitz.Matrix(1.5, 1.5)
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        return pixmap.tobytes("jpeg", jpg_quality=80)
    finally:
        if document is not None:
            document.close()
        os.remove(tmp_path)


def render_open_pdf_page_to_jpeg(document, page_index: int) -> bytes:
    """Render a page from an already opened PyMuPDF document to JPEG bytes."""
    try:
        import fitz  # PyMuPDF
    except Exception as exc:
        raise RuntimeError(
            "Hindi PDF text extraction needs PyMuPDF for rendered-page OCR. "
            "Install backend requirement PyMuPDF and redeploy."
        ) from exc

    page = document.load_page(page_index)
    matrix = fitz.Matrix(1.5, 1.5)
    pixmap = page.get_pixmap(matrix=matrix, alpha=False)
    return pixmap.tobytes("jpeg", jpg_quality=80)


def extract_rendered_pdf_page_text_with_ocr(file_bytes: bytes, page_index: int) -> str:
    """OCR one rendered PDF page when embedded text is unusable."""
    image_bytes = render_pdf_page_to_jpeg(file_bytes, page_index)
    return (extract_text_from_image_bytes(image_bytes) or "").strip()


def extract_rendered_pdf_text_with_ocr(file_bytes: bytes) -> str:
    """OCR every rendered page in a PDF when the embedded text layer is broken."""
    try:
        import fitz  # PyMuPDF
    except Exception as exc:
        raise RuntimeError(
            "Hindi PDF text extraction needs PyMuPDF for rendered-page OCR. "
            "Install backend requirement PyMuPDF and redeploy."
        ) from exc

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    document = None
    try:
        document = fitz.open(tmp_path)
        text_parts = []

        for page_index in range(document.page_count):
            image_bytes = render_open_pdf_page_to_jpeg(document, page_index)
            text_parts.append((extract_text_from_image_bytes(image_bytes) or "").strip())

        return "\n\n".join(part for part in text_parts if part).strip()
    finally:
        if document is not None:
            document.close()
        os.remove(tmp_path)


def maybe_ocr_broken_hindi_pdf_page(
    *,
    file_bytes: bytes,
    filename: str,
    page_index: int,
    text: str,
) -> str:
    """Replace one page's broken Hindi PDF text with rendered-page OCR when useful."""
    if not looks_like_broken_hindi_pdf_text(text, filename):
        return text

    ocr_text = extract_rendered_pdf_page_text_with_ocr(file_bytes, page_index)
    if is_better_hindi_extraction(ocr_text, text):
        return ocr_text

    return text


def count_words(text: str) -> int:
    """Count whitespace-separated words for OCR quality warnings."""
    return len(text.split())


def build_extracted_page(
    *,
    filename: str,
    page_number: int,
    text: str,
    extraction_method: str,
    warnings: list[str] | None = None,
) -> dict:
    """
    Build normalized page metadata for upload review.

    Low-word pages are flagged because they often indicate a blurry scan or a
    PDF page that needs OCR rather than embedded-text extraction.
    """
    clean_text = (text or "").strip()
    page_warnings = list(warnings or [])

    if count_words(clean_text) < MIN_USEFUL_PAGE_WORDS:
        page_warnings.append("Low text detected. Review OCR quality before upload.")

    return {
        "filename": filename,
        "page_number": page_number,
        "text": clean_text,
        "extraction_method": extraction_method,
        "word_count": count_words(clean_text),
        "warnings": page_warnings,
    }


def ocr_pdf_page_images(page) -> str:
    """
    OCR embedded images from one PDF page when normal PDF text is too sparse.

    Failures on individual images are ignored so one bad image does not abort the
    whole multi-page upload.
    """
    text_parts = []

    for image in getattr(page, "images", []) or []:
        try:
            image_text = extract_text_from_image_bytes(image.data)

            if image_text.strip():
                text_parts.append(image_text)
        except Exception:
            continue

    return "\n\n".join(text_parts).strip()


def extract_pages_from_pdf(filename: str, file_bytes: bytes) -> list[dict]:
    """
    Extract per-page text from a PDF, falling back to embedded-image OCR.

    Per-page output lets downstream chapter-detection group pages by
    subject/chapter after extraction.
    """
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        try:
            reader = PdfReader(tmp_path)
            pages = []

            for page_index, page in enumerate(reader.pages, start=1):
                warnings = []
                text = (page.extract_text() or "").strip()
                extraction_method = "pdf_text"

                if looks_like_broken_hindi_pdf_text(text, filename):
                    ocr_text = extract_rendered_pdf_page_text_with_ocr(
                        file_bytes,
                        page_index - 1,
                    )

                    if is_better_hindi_extraction(ocr_text, text):
                        text = ocr_text
                        extraction_method = "pdf_page_rendered_ocr"
                    else:
                        warnings.append(
                            "Hindi PDF text layer looks unreadable. Review OCR quality before upload."
                        )

                if count_words(text) < MIN_USEFUL_PAGE_WORDS:
                    ocr_text = ocr_pdf_page_images(page)

                    if ocr_text:
                        text = ocr_text
                        extraction_method = "pdf_embedded_image_ocr"
                    else:
                        warnings.append(
                            "PDF page may be scanned. No embedded image text could be extracted."
                        )

                pages.append(
                    build_extracted_page(
                        filename=filename,
                        page_number=page_index,
                        text=text,
                        extraction_method=extraction_method,
                        warnings=warnings,
                    )
                )

            return pages
        except Exception:
            return [
                build_extracted_page(
                    filename=filename,
                    page_number=page_index + 1,
                    text=maybe_ocr_broken_hindi_pdf_page(
                        file_bytes=file_bytes,
                        filename=filename,
                        page_index=page_index,
                        text=text,
                    ),
                    extraction_method="pdf_text_pymupdf",
                )
                for page_index, text in enumerate(
                    extract_pdf_text_parts_with_pymupdf(file_bytes)
                )
            ]

    finally:
        os.remove(tmp_path)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract paragraph text from a DOCX upload."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        doc = Document(tmp_path)
        return "\n".join(
            paragraph.text
            for paragraph in doc.paragraphs
            if paragraph.text.strip()
        )

    finally:
        os.remove(tmp_path)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract slide text from a PPTX upload while preserving slide order."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        prs = Presentation(tmp_path)
        text_parts = []

        for slide_number, slide in enumerate(prs.slides, start=1):
            text_parts.append(f"\nSlide {slide_number}\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)

        return "\n".join(text_parts).strip()

    finally:
        os.remove(tmp_path)


def extract_text_from_uploaded_file(filename: str, file_bytes: bytes) -> str:
    """
    Route an uploaded file to the correct text extractor based on extension.

    Images go through OCR, modern Office formats use their parsers, and legacy
    DOC/PPT files are rejected with a conversion instruction.
    """
    ext = os.path.splitext(filename.lower())[1]

    if ext in IMAGE_EXTENSIONS:
        return extract_text_from_image_bytes(file_bytes)
    
    if ext == ".txt":
        return extract_text_from_txt(file_bytes)

    if ext == ".pdf":
        return extract_text_from_pdf(file_bytes, filename)

    if ext == ".docx":
        return extract_text_from_docx(file_bytes)

    if ext == ".pptx":
        return extract_text_from_pptx(file_bytes)

    if ext in [".doc", ".ppt"]:
        raise ValueError(
            "Legacy .doc/.ppt files are not directly supported yet. "
            "Please convert them to .docx/.pptx first."
        )

    raise ValueError(
        "Unsupported file type. Supported: txt, jpg, jpeg, png, webp, pdf, docx, pptx."
    )


def extract_pages_from_uploaded_file(filename: str, file_bytes: bytes) -> list[dict]:
    """
    Return page-like extraction records for any supported upload type.

    Single-document formats are represented as one page so downstream upload
    code can process all inputs through one grouping pipeline.
    """
    ext = os.path.splitext(filename.lower())[1]

    if ext in IMAGE_EXTENSIONS:
        return [
            build_extracted_page(
                filename=filename,
                page_number=1,
                text=extract_text_from_image_bytes(file_bytes),
                extraction_method="image_ocr",
            )
        ]

    if ext == ".pdf":
        return extract_pages_from_pdf(filename, file_bytes)

    if ext == ".txt":
        return [
            build_extracted_page(
                filename=filename,
                page_number=1,
                text=extract_text_from_txt(file_bytes),
                extraction_method="text",
            )
        ]

    if ext == ".docx":
        return [
            build_extracted_page(
                filename=filename,
                page_number=1,
                text=extract_text_from_docx(file_bytes),
                extraction_method="docx_text",
            )
        ]

    if ext == ".pptx":
        return [
            build_extracted_page(
                filename=filename,
                page_number=1,
                text=extract_text_from_pptx(file_bytes),
                extraction_method="pptx_text",
            )
        ]

    if ext in [".doc", ".ppt"]:
        raise ValueError(
            "Legacy .doc/.ppt files are not directly supported yet. "
            "Please convert them to .docx/.pptx first."
        )

    raise ValueError(
        "Unsupported file type. Supported: txt, jpg, jpeg, png, webp, pdf, docx, pptx."
    )
